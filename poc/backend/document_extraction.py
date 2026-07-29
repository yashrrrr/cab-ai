"""
Supporting Document Extraction — text extraction + LLM field auto-fill
Powers the "upload a BRD/FRD/PRD/RFC document" flow on RFC submission: pulls
raw text out of the uploaded file, then asks the model to map that text onto
RFC form fields.

Supported formats: PDF, Word (.docx), PowerPoint (.pptx), Excel (.xlsx).
Legacy binary formats (.doc, .ppt, .xls) are NOT supported — parsing those
reliably needs much heavier tooling (no good pure-Python reader); ask
requestors to save as the modern Office Open XML format instead.
"""

from openai import OpenAI
from typing import Dict
import json
import os

from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
import openpyxl

# Same provider/client as cab_orchestrator.py — GitHub Models endpoint via the
# OpenAI SDK. Kept as its own small client here rather than importing from
# cab_orchestrator to avoid coupling document extraction to CAB session code.
api_key = os.getenv("OPENAI_API_KEY")
if not api_key:
    raise ValueError("OPENAI_API_KEY environment variable not set")

client = OpenAI(
    api_key=api_key,
    base_url="https://models.inference.ai.azure.com"
)

# Max characters of extracted document text sent to the field-extraction
# prompt. Keeps the request within a reasonable token budget for long PDFs.
MAX_EXTRACTION_CHARS = 12000

ALLOWED_FIELDS = {
    "title",
    "description",
    "business_justification",
    "implementation_plan",
    "test_cases",
    "back_out_plan",
    "affected_systems",
}

FIELD_EXTRACTION_SYSTEM_PROMPT = """You extract structured change-request fields from a PRD/BRD/RFC document.
Read the document text and identify content for each field below, if present. Do not invent information that
isn't in the document — omit a field entirely rather than guess."""

FIELD_EXTRACTION_PROMPT_TEMPLATE = """Read this document and extract values for as many of the following RFC fields as the document actually supports:

- title: short change title
- description: what the change is
- business_justification: why this change is needed / business value
- implementation_plan: how the change will be carried out
- test_cases: testing performed or planned
- back_out_plan: rollback/recovery procedure
- affected_systems: list of affected systems/services

Output ONLY a JSON object (no prose, no markdown fences) with whatever subset of these keys the document supports.
"affected_systems" should be a JSON array of strings if present. Omit any field not clearly supported by the text.

DOCUMENT TEXT:
{document_text}
"""


def extract_pdf_text(file_path: str) -> str:
    """
    Extract all text from a PDF file. Raises on unreadable/corrupt PDFs —
    callers should surface that as a 400 to the requestor.
    """
    reader = PdfReader(file_path)
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n".join(pages).strip()


def extract_docx_text(file_path: str) -> str:
    """Extract paragraph and table text from a Word (.docx) file."""
    doc = DocxDocument(file_path)
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    parts.append(cell.text)
    return "\n".join(parts).strip()


def extract_pptx_text(file_path: str) -> str:
    """Extract text from every text frame on every slide of a PowerPoint (.pptx) file."""
    prs = Presentation(file_path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = "\n".join(p.text for p in shape.text_frame.paragraphs if p.text.strip())
                if text:
                    parts.append(text)
    return "\n".join(parts).strip()


def extract_xlsx_text(file_path: str) -> str:
    """Extract cell values (row by row, sheet by sheet) from an Excel (.xlsx) file."""
    workbook = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
    parts = []
    for sheet in workbook.worksheets:
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None and str(c).strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts).strip()


EXTRACTORS_BY_EXTENSION = {
    ".pdf": extract_pdf_text,
    ".docx": extract_docx_text,
    ".pptx": extract_pptx_text,
    ".xlsx": extract_xlsx_text,
}


def extract_document_text(file_path: str, filename: str) -> str:
    """
    Dispatch to the right extractor based on the file's extension.
    Raises ValueError for unsupported extensions, and whatever the
    underlying library raises for unreadable/corrupt files — callers
    should surface both as a 400 to the requestor.
    """
    ext = os.path.splitext(filename.lower())[1]
    extractor = EXTRACTORS_BY_EXTENSION.get(ext)
    if extractor is None:
        raise ValueError(f"Unsupported file type: {ext or filename}")
    return extractor(file_path)


def extract_rfc_fields(document_text: str) -> Dict:
    """
    Ask the LLM to map document text onto RFC submission fields.

    Best-effort only: any failure (API error, malformed/non-JSON response,
    empty extraction) yields an empty dict rather than raising, since a
    weak/failed auto-fill should never block the document from being
    attached — the requestor just fills the form manually.
    """
    if not document_text.strip():
        return {}

    prompt = FIELD_EXTRACTION_PROMPT_TEMPLATE.format(
        document_text=document_text[:MAX_EXTRACTION_CHARS]
    )

    try:
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": FIELD_EXTRACTION_SYSTEM_PROMPT},
                {"role": "user", "content": prompt},
            ],
            max_tokens=1200,
            temperature=0,
        )
        raw = response.choices[0].message.content or ""
        start = raw.find("{")
        if start == -1:
            return {}
        parsed, _ = json.JSONDecoder().raw_decode(raw[start:])
        if not isinstance(parsed, dict):
            return {}

        result = {}
        for key, value in parsed.items():
            if key not in ALLOWED_FIELDS:
                continue
            if key == "affected_systems":
                if isinstance(value, list):
                    systems = [str(s).strip() for s in value if str(s).strip()]
                elif isinstance(value, str) and value.strip():
                    systems = [s.strip() for s in value.split(",") if s.strip()]
                else:
                    continue
                if systems:
                    result[key] = systems
            elif isinstance(value, str) and value.strip():
                result[key] = value.strip()

        return result
    except Exception:
        return {}
